"""Generate 11-slide advisor presentation as .pptx.

Narrative arc:
  1  Title
  2  Motivation — two kinds of DRAM errors
  3  Why BCH fails on burst (theory diagram)
  4  BCH burst failure (empirical)
  5  Our approach — Feistel + CRC DAG
  6  Burst = Random (empirical proof)
  7  RS comparison (empirical)
  8  Overhead zones — all schemes vs block size
  9  Correction capability
  10 Contribution summary table
  11 Conclusion
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Dimensions (16:9 widescreen) ─────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
GREEN  = RGBColor(0x1E, 0x8B, 0x4C)
RED    = RGBColor(0xC0, 0x39, 0x2B)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
MGRAY  = RGBColor(0x88, 0x88, 0x88)
LNAVY  = RGBColor(0xAA, 0xCC, 0xEE)
DNAVY  = RGBColor(0x10, 0x1E, 0x35)

R = os.path.join(os.path.dirname(__file__), "results")
FIG = {
    "bch_burst":      os.path.join(R, "fig_bch_burst",    "fig_bch_burst.png"),
    "fig4":           os.path.join(R, "fig4",             "fig4_burst_resilience.png"),
    "rs_empirical":   os.path.join(R, "fig_rs_empirical", "fig_rs_empirical.png"),
    "overhead_zones": os.path.join(R, "fig_overhead_zones","fig_overhead_zones.png"),
    "fig1":           os.path.join(R, "fig1",             "fig1_headline.png"),
    "fig2":           os.path.join(R, "fig2",             "fig2_overhead_comparison.png"),
}

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_bullets(slide, items, x, y, w, h, size=16, color=DGRAY, bullet="▸", spacing=6):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txb


def add_table(slide, headers, rows, x, y, w, h,
              hdr_fill=NAVY, hdr_fg=WHITE, font_size=15):
    n_cols, n_rows = len(headers), len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    col_w = w // n_cols
    for c in range(n_cols):
        tbl.columns[c].width = col_w

    def cell(r, c, text, fill, fg, bold=False, align=PP_ALIGN.CENTER):
        cl = tbl.cell(r, c)
        cl.text = text
        cl.fill.solid()
        cl.fill.fore_color.rgb = fill
        p = cl.text_frame.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = fg

    for c, h_label in enumerate(headers):
        cell(0, c, h_label, hdr_fill, hdr_fg, bold=True)
    for r, row in enumerate(rows):
        bg = LGRAY if r % 2 == 1 else WHITE
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            # Highlight our scheme row in light blue
            if row[0].startswith("Ours"):
                bg = RGBColor(0xE8, 0xF4, 0xFF)
            cell(r + 1, c, val, bg, DGRAY, bold=row[0].startswith("Ours"), align=align)


def embed(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        add_rect(slide, x, y, w, h, fill=LGRAY, line=MGRAY)
        add_text(slide, f"[pending: {os.path.basename(path)}]",
                 x, y + h // 2 - Inches(0.3), w, Inches(0.5),
                 size=13, color=MGRAY, align=PP_ALIGN.CENTER)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def slide_header(slide, title, subtitle=None, bg=WHITE):
    add_rect(slide, 0, 0, W, H, fill=bg)
    add_rect(slide, 0, 0, W, Inches(1.1), fill=NAVY)
    add_text(slide, title,
             Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.72),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.76), Inches(12.6), Inches(0.36),
                 size=14, color=LNAVY)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.04), fill=ACCENT)


def caption(slide, text):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill=RGBColor(0xF0, 0xF4, 0xF8))
    add_text(slide, text, Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.36),
             size=13, color=DGRAY, italic=True)


def full_fig(slide, key, caption_text):
    """Embed a figure that fills the content area with a caption bar at bottom."""
    embed(slide, FIG[key], Inches(0.2), Inches(1.2), Inches(12.9), Inches(5.8))
    caption(slide, caption_text)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(2.6), W, Inches(2.6), fill=DNAVY)
add_text(sl, "Approximate ECC via\nFeistel-Permuted CRC Hashing",
         Inches(1), Inches(1.3), Inches(11.3), Inches(2.1),
         size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Burst resilience and overhead scaling for AI accelerator memory",
         Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
         size=20, color=LNAVY, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(5.5), Inches(4.0), Inches(2.33), Inches(0.05), fill=ACCENT)
add_text(sl, "Research Presentation  ·  2026",
         Inches(1), Inches(4.2), Inches(11.3), Inches(0.4),
         size=16, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)
note(sl, "Lead with burst immunity — that is the primary contribution.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Motivation: Two Kinds of DRAM Errors",
             "AI accelerators store model weights in HBM/DRAM — soft errors corrupt inference silently")

# Left column — text
add_bullets(sl, [
    "Modern AI accelerators (A100, H100, TPU) keep billions of weight parameters in DRAM",
    "Soft errors occur at rates of 10⁻⁶–10⁻³ per bit per hour — high enough to matter at scale",
    "Two physically distinct failure modes require different protection:",
], Inches(0.4), Inches(1.3), Inches(6.0), Inches(2.2), size=15, spacing=8)

# Two error type boxes
for i, (label, desc, color, bits) in enumerate([
    ("Random errors", "Cosmic rays / alpha particles\nIsolated, uniformly scattered bits", ACCENT,
     "0 1 0 1 1 0 0 1 0 0 1 0\n0 0 1 0 0 1 0 0 0 1 0 0\n1 0 0 0 1 0 0 1 0 0 0 1"),
    ("Burst errors",  "Row-hammer / retention failures\nContiguous run of corrupted bits", RED,
     "0 1 0 1 1 0 0 1 0 0 1 0\n0 0 ■ ■ ■ ■ ■ ■ 0 1 0 0\n1 0 0 0 1 0 0 1 0 0 0 1"),
]):
    bx = Inches(0.4 + i * 6.0)
    by = Inches(3.6)
    add_rect(sl, bx, by, Inches(5.7), Inches(3.4), fill=LGRAY, line=color)
    add_rect(sl, bx, by, Inches(5.7), Inches(0.45), fill=color)
    add_text(sl, label, bx + Inches(0.15), by + Inches(0.06),
             Inches(5.4), Inches(0.35), size=16, bold=True, color=WHITE)
    add_text(sl, desc, bx + Inches(0.2), by + Inches(0.55),
             Inches(5.3), Inches(0.7), size=13, color=DGRAY)
    add_text(sl, bits, bx + Inches(0.2), by + Inches(1.35),
             Inches(5.3), Inches(1.5), size=13, color=DGRAY,
             wrap=False)

add_rect(sl, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.45),
         fill=RGBColor(0xFF, 0xF3, 0xCD))
add_text(sl, "No single existing code handles both cheaply — you must choose one or deploy two separate systems.",
         Inches(0.6), Inches(6.93), Inches(12.1), Inches(0.4),
         size=14, bold=True, color=RGBColor(0x7D, 0x5A, 0x00))
note(sl, "Motivate the need — existing ECC is designed for one model. BCH assumes BSC (random). RS is symbol-based (burst). Neither is free.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why BCH fails on burst (theory)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Why BCH Fails on Burst Errors",
             "BCH is designed for the Binary Symmetric Channel — independent random errors")

# Draw 16 BCH blocks as a row of rectangles
n_blocks = 16
bw = Inches(0.72)
bh = Inches(0.9)
bx0 = Inches(0.35)
by  = Inches(2.1)
gap = Inches(0.04)

for i in range(n_blocks):
    bx = bx0 + i * (bw + gap)
    add_rect(sl, bx, by, bw, bh, fill=RGBColor(0xD6, 0xE4, 0xF0), line=ACCENT)
    add_text(sl, f"B{i}", bx, by + Inches(0.05), bw, Inches(0.3),
             size=9, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, "256 bits\nt=19", bx, by + Inches(0.36), bw, Inches(0.5),
             size=8, color=DGRAY, align=PP_ALIGN.CENTER)

add_text(sl, "4096-bit data block  (16 × 256-bit BCH sub-blocks, each correcting t = 19 random bits)",
         Inches(0.35), Inches(1.75), Inches(12.6), Inches(0.32),
         size=13, color=DGRAY)

# Burst bar spanning blocks 3-6 (20 bits)
burst_start = 3
burst_span  = 3
burst_x = bx0 + burst_start * (bw + gap)
burst_w = burst_span * (bw + gap) + Inches(0.08)
add_rect(sl, burst_x, by - Inches(0.55), burst_w, Inches(0.4),
         fill=RGBColor(0xFF, 0xCC, 0xCC), line=RED)
add_text(sl, "20-bit burst", burst_x, by - Inches(0.56),
         burst_w, Inches(0.38), size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
# Arrow down
add_rect(sl, burst_x + burst_w / 2 - Inches(0.03),
         by - Inches(0.14), Inches(0.06), Inches(0.14), fill=RED)

# Highlight the hit block with red
hit_bx = bx0 + burst_start * (bw + gap)
add_rect(sl, hit_bx, by, bw * 1.5 + gap, bh, fill=RGBColor(0xFF, 0xE5, 0xE5), line=RED)
add_text(sl, "✗ FAIL\n>19 bits\nin block",
         hit_bx + Inches(0.02), by + Inches(0.05), bw * 1.5, bh - Inches(0.1),
         size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Explanation text
add_rect(sl, Inches(0.35), Inches(3.2), Inches(12.6), Inches(0.06), fill=MGRAY)
add_bullets(sl, [
    "A 20-bit burst hits B3–B4, putting 20 errors into block B3 alone  (t = 19 → uncorrectable)",
    "BCH has no way to redistribute the burst — each block decodes independently",
    "BCH corrects t random bits per block; burst errors concentrate in one block and exceed t immediately",
    "Interleaving can help — but requires reordering all data before write and after read (latency cost)",
], Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.6), size=15, spacing=10)

note(sl, "Key: BCH's independence assumption is violated by burst. Each block decodes alone.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BCH burst failure (empirical)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "BCH Burst Failure — Empirical Confirmation",
             "BCH(256-bit blocks, t=13) vs our scheme — success rate vs contiguous burst length")
full_fig(sl, "bch_burst",
         "BCH collapses the moment burst length exceeds t per block. Our scheme stays at 100% success across all burst lengths tested.")
note(sl, "The cliff at t=13 is sharp — this is a hard mathematical boundary, not gradual degradation.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Our approach
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Our Approach: Feistel-Permuted CRC Hash DAG",
             "A single scheme that handles both error types via uniform bit redistribution")

# Left: pipeline boxes
steps = [
    ("①  Feistel Permutation",    "Map source bits onto an N×N grid\nvia a keyed Feistel shuffle"),
    ("②  CRC Hash Nodes",         "Compute CRC-8/16/32 over each\nrow and column group"),
    ("③  Hash DAG",                "Build constraint graph: edges\nbetween overlapping nodes"),
    ("④  DAG-Guided Solver",       "Enumerate flip combinations\nuntil all hashes agree"),
]
box_h = Inches(1.15)
box_gap = Inches(0.18)
box_y0 = Inches(1.25)
for i, (title, desc) in enumerate(steps):
    bx = Inches(0.4)
    by = box_y0 + i * (box_h + box_gap)
    add_rect(sl, bx, by, Inches(6.0), box_h,
             fill=RGBColor(0xEA, 0xF4, 0xFF), line=ACCENT)
    add_rect(sl, bx, by, Inches(6.0), Inches(0.38), fill=ACCENT)
    add_text(sl, title, bx + Inches(0.12), by + Inches(0.04),
             Inches(5.8), Inches(0.32), size=14, bold=True, color=WHITE)
    add_text(sl, desc, bx + Inches(0.15), by + Inches(0.45),
             Inches(5.7), Inches(0.65), size=13, color=DGRAY)
    # Arrow to next step
    if i < len(steps) - 1:
        ax = bx + Inches(2.8)
        ay = by + box_h
        add_rect(sl, ax, ay, Inches(0.06), box_gap, fill=MGRAY)

# Right: key insight box
rx = Inches(7.0)
add_rect(sl, rx, Inches(1.25), Inches(5.95), Inches(5.85),
         fill=RGBColor(0xF0, 0xFF, 0xF0), line=GREEN)
add_rect(sl, rx, Inches(1.25), Inches(5.95), Inches(0.48), fill=GREEN)
add_text(sl, "Key Insight", rx + Inches(0.15), Inches(1.29),
         Inches(5.6), Inches(0.38), size=16, bold=True, color=WHITE)

add_text(sl,
    "The Feistel permutation scatters contiguous "
    "burst bits uniformly across the 2D grid.\n\n"
    "From the solver's perspective, burst errors "
    "are indistinguishable from random errors.\n\n"
    "This means one scheme — one overhead budget — "
    "corrects both error patterns.",
    rx + Inches(0.2), Inches(1.85), Inches(5.55), Inches(3.5),
    size=14, color=DGRAY, wrap=True)

add_rect(sl, rx, Inches(5.6), Inches(5.95), Inches(1.5),
         fill=RGBColor(0xE8, 0xF4, 0xFF), line=ACCENT)
add_text(sl, "Overhead = 2 × HASH_BITS / √L\n"
             "BER-independent — set once at deployment",
         rx + Inches(0.2), Inches(5.65), Inches(5.6), Inches(1.3),
         size=14, color=NAVY, bold=False)

note(sl, "The Feistel permutation is the mechanism. Without it, burst errors cluster in one row/column and overwhelm that node.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Burst = Random (empirical proof)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Empirical Proof: Burst Errors = Random Errors",
             "4096-bit block, CRC-32, 100% overhead — success rate and solver effort vs injected flip count")
full_fig(sl, "fig4",
         "Both curves overlap exactly across all flip counts. The Feistel permutation eliminates any distinction between burst and random for our solver.")
note(sl, "This is the central empirical result. The overlap is exact — not approximate.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RS comparison (empirical)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Reed-Solomon Comparison: Burst vs Random",
             "RS over GF(2⁸), 255-byte codewords — empirical success rate at two overhead levels")
full_fig(sl, "rs_empirical",
         "RS at 38% overhead handles burst but collapses on random. At 125% overhead RS handles both — but costs more than our 100% overhead scheme.")
note(sl, "RS is symbol-based: cheap for burst (contiguous bits hit few symbols) but expensive for random (39% symbol error rate at 6% BER).")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Overhead zones
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Overhead vs Block Size — All Schemes, All BER Levels",
             "Left: RS vs Ours  ·  Right: BCH vs Ours  ·  Each BER level shown (lighter = lower BER)")
full_fig(sl, "overhead_zones",
         "Our O(1/√L) overhead beats RS (for equal correction) above ~7K bits, and BCH above ~15K bits. Crossover shifts left at higher BER.")
note(sl, "BCH is flat because it tiles at the natural 256-bit codeword. RS starts high because small data underfills the larger 2040-bit natural codeword.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Correction capability
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Correction Capability vs Flip Count",
             "4096-bit block — success rate, solve time, and hash checks for CRC-8/16/32")
full_fig(sl, "fig1",
         "CRC-32 (100% overhead) corrects 200+ bit flips at 100% success. CRC-16 (50%) handles ~100 flips. CRC-8 (25%) handles ~60 flips.")
note(sl, "Show the three operating points. The 32-bit config is the comparison point for RS and BCH slides.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Contribution summary table
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "Contribution Summary",
             "Single scheme that handles both error types with sub-linear overhead growth")

headers = ["Scheme", "Burst errors", "Random errors", "Overhead at 6% BER", "Scales O(1/√L)"]
rows = [
    ["BCH",    "✗  Fails",          "✓  53% overhead",    "53%  (flat)",             "✗"],
    ["RS",     "✓  38% overhead",   "✓  125% overhead",   "38–125%  (flat)",          "✗"],
    ["Ours ✦", "✓  100% overhead",  "✓  100% overhead",   "100% → shrinks with L",   "✓"],
]
add_table(sl, headers, rows,
          Inches(0.4), Inches(1.35), Inches(12.5), Inches(2.6),
          font_size=15)

# Takeaway bullets below the table
add_bullets(sl, [
    "BCH is cheaper for random at small blocks — but completely blind to burst errors",
    "RS handles burst cheaply — but random errors cost 3× more due to symbol granularity",
    "Our scheme pays the same overhead for both, and that overhead shrinks as L grows",
    "Above ~7K bits: our scheme is cheaper than RS for equal correction.  Above ~15K bits: cheaper than BCH.",
], Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.0), size=15, spacing=9)

note(sl, "Emphasize: BCH and RS each solve half the problem. We solve both halves with one mechanism.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(0.85), W, Inches(0.05), fill=ACCENT)

add_text(sl, "Conclusions", Inches(0.6), Inches(0.15), Inches(12.0), Inches(0.65),
         size=32, bold=True, color=WHITE)

contributions = [
    ("Feistel equalization",
     "The Feistel permutation distributes burst bits uniformly — empirically confirmed. "
     "Burst and random errors are identical for our solver."),
    ("O(1/√L) overhead scaling",
     "Overhead formula: 2·HASH_BITS / √L. BER-independent. "
     "Beats RS above ~7K bits, BCH above ~15K bits — and keeps shrinking."),
    ("Single unified scheme",
     "One deployment replaces two separate ECC systems. "
     "No interleaving, no field switching, no mode selection."),
]

for i, (title, body) in enumerate(contributions):
    cy = Inches(1.15) + i * Inches(1.95)
    add_rect(sl, Inches(0.5), cy, Inches(12.3), Inches(1.75),
             fill=RGBColor(0x0F, 0x1D, 0x33), line=ACCENT)
    add_rect(sl, Inches(0.5), cy, Inches(0.28), Inches(1.75), fill=ACCENT)
    add_text(sl, f"  {title}", Inches(0.85), cy + Inches(0.12),
             Inches(11.8), Inches(0.45), size=17, bold=True, color=LNAVY)
    add_text(sl, body, Inches(0.85), cy + Inches(0.62),
             Inches(11.8), Inches(1.0), size=14, color=WHITE, wrap=True)

add_text(sl, "Thank you", Inches(0.6), Inches(6.95), Inches(12.0), Inches(0.45),
         size=15, color=MGRAY, italic=True)

note(sl, "End here. Backup slides available if asked about false-positive probability, solver complexity, or multi-burst patterns.")


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(__file__), "advisor_slides.pptx")
prs.save(OUT)
print(f"Saved: {OUT}  ({prs.slides.__len__()} slides)")
